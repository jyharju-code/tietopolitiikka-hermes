"""Local Telegram ingestion for Tietopolitiikka Hermes.

Every event accepted from the exact configured Telegram supergroup is written
to a durable local spool before agent routing. Public URLs and cached Telegram
attachments are extracted locally and written through OpenViking content/write.
"""

from __future__ import annotations

import asyncio
import hashlib
import html
import ipaddress
import json
import mimetypes
import os
import re
import shutil
import socket
import tempfile
from datetime import datetime, timezone
from pathlib import Path
from typing import Any
from urllib.parse import urljoin, urlparse


DATA_ROOT = Path(os.environ.get("HERMES_HOME", "/opt/data"))
SPOOL_ROOT = DATA_ROOT / "ingest-spool"
FILE_ROOT = DATA_ROOT / "ingest-files"
MAX_DOWNLOAD_BYTES = 12 * 1024 * 1024
MAX_EXTRACTED_CHARS = 2_000_000
CHUNK_CHARS = 12_000
MAX_SITE_DOCUMENTS = 100
HTTP_URL_PATTERN = re.compile(r"https?://[^\s<>\]\[(){}\"']+", re.IGNORECASE)
BARE_DOMAIN_PATTERN = re.compile(
    r"(?<![A-Za-z0-9@:/])(?:[A-Za-z0-9](?:[A-Za-z0-9-]{0,61}[A-Za-z0-9])?\.)+"
    r"[A-Za-z]{2,63}(?:/[^\s<>\]\[(){}\"']*)?",
    re.IGNORECASE,
)
_WORKER_TASK: asyncio.Task[None] | None = None
_WORKER_WAKEUP: asyncio.Event | None = None


def _clean_url(value: str) -> str:
    return value.rstrip(".,;:!?'")


def _extract_urls(text: str) -> list[str]:
    """Extract explicit and bare web addresses and normalize bare domains."""
    matches: list[tuple[int, int, str]] = [
        (match.start(), match.end(), _clean_url(match.group(0)))
        for match in HTTP_URL_PATTERN.finditer(text)
    ]
    occupied = [(start, end) for start, end, _value in matches]
    for match in BARE_DOMAIN_PATTERN.finditer(text):
        if any(match.start() < end and match.end() > start for start, end in occupied):
            continue
        matches.append((match.start(), match.end(), f"https://{_clean_url(match.group(0))}"))
    return sorted({value for _start, _end, value in matches})


def _requests_site_document_crawl(text: str) -> bool:
    compact = re.sub(r"\s+", " ", text.lower())
    return bool(
        re.search(r"\b(?:kaikki|all)\b.{0,60}\b(?:pdf|dokument)", compact)
        or re.search(r"\b(?:pdf|dokument)\w*\b.{0,60}\b(?:kaikki|all)\b", compact)
    )


def _sha(value: str | bytes) -> str:
    if isinstance(value, str):
        value = value.encode("utf-8", errors="replace")
    return hashlib.sha256(value).hexdigest()


def _file_sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _safe_segment(value: Any, fallback: str) -> str:
    text = re.sub(r"[^A-Za-z0-9._-]+", "-", str(value or "").strip()).strip("-.")
    return (text[:80] or fallback)


def _timestamp(data: dict[str, Any]) -> datetime:
    raw = data.get("timestamp") or data.get("messageTimestamp")
    try:
        number = float(raw)
        if number > 10_000_000_000:
            number /= 1000
        return datetime.fromtimestamp(number, tz=timezone.utc)
    except (TypeError, ValueError, OSError):
        return datetime.now(tz=timezone.utc)


def _event_value(event: Any, name: str, default: Any = "") -> Any:
    value = getattr(event, name, default)
    return default if value is None else value


def _telegram_event_data(event: Any) -> dict[str, Any]:
    source = getattr(event, "source", None)
    metadata = getattr(event, "metadata", None)
    if not isinstance(metadata, dict):
        metadata = {}
    return {
        "chatId": str(getattr(source, "chat_id", "") or metadata.get("chat_id") or ""),
        "chatName": str(metadata.get("chat_title") or metadata.get("chat_name") or ""),
        "senderId": str(getattr(source, "user_id", "") or metadata.get("user_id") or ""),
        "senderName": str(metadata.get("sender_name") or metadata.get("from_name") or ""),
        "messageId": str(_event_value(event, "message_id", metadata.get("message_id", ""))),
        "timestamp": metadata.get("timestamp") or metadata.get("date"),
        "threadId": str(getattr(source, "thread_id", "") or metadata.get("message_thread_id") or ""),
        "isGroup": True,
    }


def is_allowed_telegram_event(event: Any) -> bool:
    """Allow ingestion only from the exact configured Telegram supergroup."""
    allowed_chat = os.environ.get("TELEGRAM_GROUP_ID", "").strip()
    if not allowed_chat:
        return False
    data = _telegram_event_data(event)
    return data["chatId"] == allowed_chat


def _copy_media(event: Any, message_key: str) -> list[dict[str, str]]:
    FILE_ROOT.mkdir(parents=True, exist_ok=True)
    copied: list[dict[str, str]] = []
    media_urls = list(getattr(event, "media_urls", []) or [])
    media_types = list(getattr(event, "media_types", []) or [])
    for index, value in enumerate(media_urls):
        media_type = str(media_types[index] if index < len(media_types) else "")
        path = Path(str(value))
        if not path.is_absolute() or not path.is_file():
            copied.append({"source": str(value), "path": "", "mime": media_type})
            continue
        suffix = path.suffix.lower()[:16]
        destination = FILE_ROOT / f"{message_key}-{index:02d}{suffix}"
        if not destination.exists():
            shutil.copy2(path, destination)
        copied.append({"source": str(value), "path": str(destination), "mime": media_type})
    return copied


def _create_spool(event: Any, data: dict[str, Any]) -> Path:
    chat_id = str(data.get("chatId") or "unknown-chat")
    message_id = str(data.get("messageId") or "")
    body = str(getattr(event, "text", "") or data.get("body") or "")
    created = _timestamp(data)
    key = _sha(f"{chat_id}\0{message_id or body}\0{created.isoformat()}")[:32]
    urls = _extract_urls(body)
    payload = {
        "version": 1,
        "key": key,
        "platform": "telegram",
        "chat_id": chat_id,
        "chat_name": str(data.get("chatName") or ""),
        "sender_id": str(data.get("senderId") or data.get("from") or ""),
        "sender_name": str(data.get("senderName") or ""),
        "message_id": message_id,
        "thread_id": str(data.get("threadId") or ""),
        "timestamp": created.isoformat(),
        "body": body,
        "urls": urls,
        "crawl_site_documents": _requests_site_document_crawl(body),
        "media": _copy_media(event, key),
    }
    SPOOL_ROOT.mkdir(parents=True, exist_ok=True)
    target = SPOOL_ROOT / f"{key}.json"
    if not target.exists():
        temporary = target.with_suffix(".tmp")
        temporary.write_text(json.dumps(payload, ensure_ascii=False), encoding="utf-8")
        os.chmod(temporary, 0o600)
        temporary.replace(target)
    return target


def _headers() -> dict[str, str]:
    api_key = os.environ.get("OPENVIKING_API_KEY", "")
    return {
        "Content-Type": "application/json",
        "X-API-Key": api_key,
        "Authorization": f"Bearer {api_key}",
        "X-OpenViking-Account": os.environ.get("OPENVIKING_ACCOUNT", "tietopolitiikka"),
        "X-OpenViking-User": os.environ.get("OPENVIKING_USER", "telegram-core"),
        "X-OpenViking-Actor-Peer": os.environ.get("OPENVIKING_AGENT", "tietopolitiikka-hermes"),
    }


async def _write_content(uri: str, content: str) -> None:
    import httpx

    endpoint = os.environ.get("OPENVIKING_ENDPOINT", "http://openviking:1933").rstrip("/")
    async with httpx.AsyncClient(timeout=30.0) as client:
        response = await client.post(
            f"{endpoint}/api/v1/content/write",
            headers=_headers(),
            json={"uri": uri, "content": content, "mode": "create"},
        )
    if response.status_code == 409:
        return
    response.raise_for_status()
    result = response.json()
    if isinstance(result, dict) and result.get("status") == "error":
        raise RuntimeError(f"OpenViking write failed for {uri}")


def _chunk_text(text: str) -> list[str]:
    text = text[:MAX_EXTRACTED_CHARS]
    if not text:
        return [""]
    chunks: list[str] = []
    position = 0
    while position < len(text):
        end = min(position + CHUNK_CHARS, len(text))
        if end < len(text):
            boundary = text.rfind("\n", position, end)
            if boundary > position + CHUNK_CHARS // 2:
                end = boundary
        chunks.append(text[position:end].strip())
        position = end
    return chunks


async def _write_chunks(base_uri: str, heading: str, text: str) -> None:
    chunks = _chunk_text(text)
    total = len(chunks)
    for index, chunk in enumerate(chunks, start=1):
        suffix = "" if total == 1 else f"-part-{index:04d}"
        content = f"# {heading}\n\nPart: {index}/{total}\n\n{chunk}\n"
        await _write_content(f"{base_uri}{suffix}.md", content)


def _message_markdown(payload: dict[str, Any]) -> str:
    attachment_lines = [
        f"- {item.get('path') or item.get('source')} ({item.get('mime') or 'unknown'})"
        for item in payload.get("media", [])
    ]
    url_lines = [f"- {url}" for url in payload.get("urls", [])]
    return "\n".join(
        [
            "# Telegram message",
            "",
            f"Chat: {payload.get('chat_name') or payload.get('chat_id')}",
            f"Sender: {payload.get('sender_name') or payload.get('sender_id')}",
            f"Timestamp: {payload.get('timestamp')}",
            f"Message ID: {payload.get('message_id')}",
            f"Topic ID: {payload.get('thread_id') or 'general'}",
            "",
            "## Message",
            "",
            payload.get("body") or "[No text body]",
            "",
            "## URLs",
            "",
            *(url_lines or ["[None]"]),
            "",
            "## Attachments",
            "",
            *(attachment_lines or ["[None]"]),
        ]
    )


async def _validate_public_host(hostname: str) -> None:
    loop = asyncio.get_running_loop()
    records = await loop.getaddrinfo(hostname, None, type=socket.SOCK_STREAM)
    if not records:
        raise ValueError("URL host did not resolve")
    for record in records:
        address = ipaddress.ip_address(record[4][0])
        if not address.is_global:
            raise ValueError("Private or non-global URL address is blocked")


async def _download_public_url(url: str) -> tuple[bytes, str, str]:
    import httpx

    current = url
    async with httpx.AsyncClient(timeout=20.0, follow_redirects=False) as client:
        for _ in range(5):
            parsed = urlparse(current)
            if parsed.scheme not in {"http", "https"} or not parsed.hostname:
                raise ValueError("Only public HTTP and HTTPS URLs are supported")
            await _validate_public_host(parsed.hostname)
            async with client.stream("GET", current, headers={"User-Agent": "TietopolitiikkaHermes/1.0"}) as response:
                if response.status_code in {301, 302, 303, 307, 308}:
                    location = response.headers.get("location")
                    if not location:
                        raise ValueError("Redirect without location")
                    current = urljoin(current, location)
                    continue
                response.raise_for_status()
                content_type = response.headers.get("content-type", "application/octet-stream").split(";", 1)[0]
                chunks: list[bytes] = []
                size = 0
                async for chunk in response.aiter_bytes():
                    size += len(chunk)
                    if size > MAX_DOWNLOAD_BYTES:
                        raise ValueError("URL content exceeds the local ingestion limit")
                    chunks.append(chunk)
                return b"".join(chunks), content_type, current
    raise ValueError("Too many redirects")


def _html_text(content: bytes) -> str:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(content, "html.parser")
    for node in soup(["script", "style", "noscript", "svg"]):
        node.decompose()
    title = soup.title.get_text(" ", strip=True) if soup.title else ""
    body = soup.get_text("\n", strip=True)
    return f"{title}\n\n{body}".strip()


def _extract_path(path: Path, mime_type: str = "") -> str:
    suffix = path.suffix.lower()
    if suffix in {".txt", ".md", ".csv", ".json", ".xml", ".yaml", ".yml", ".log", ".py", ".js", ".ts", ".css"}:
        return path.read_text(encoding="utf-8", errors="replace")
    if suffix in {".html", ".htm"}:
        return _html_text(path.read_bytes())
    if suffix == ".pdf":
        from pypdf import PdfReader

        return "\n\n".join(page.extract_text() or "" for page in PdfReader(str(path)).pages)
    if suffix == ".docx":
        from docx import Document

        document = Document(str(path))
        return "\n".join(paragraph.text for paragraph in document.paragraphs)
    if suffix == ".pptx":
        from pptx import Presentation

        presentation = Presentation(str(path))
        return "\n".join(
            shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")
        )
    if suffix in {".xlsx", ".xlsm"}:
        from openpyxl import load_workbook

        workbook = load_workbook(str(path), read_only=True, data_only=True)
        rows: list[str] = []
        for sheet in workbook.worksheets:
            rows.append(f"Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                rows.append("\t".join("" if value is None else str(value) for value in row))
        return "\n".join(rows)
    if mime_type.startswith("image/") or suffix in {".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff", ".bmp"}:
        import pytesseract
        from PIL import Image

        with Image.open(path) as image:
            return pytesseract.image_to_string(image, lang="fin+eng")
    try:
        return path.read_text(encoding="utf-8", errors="strict")
    except (UnicodeDecodeError, OSError):
        size = path.stat().st_size
        digest = _sha(path.read_bytes())
        return f"Binary attachment archived locally. File: {path.name}\nMIME: {mime_type}\nBytes: {size}\nSHA-256: {digest}"


async def _index_url(url: str) -> None:
    url_key = _sha(url)[:32]
    base_uri = f"viking://user/telegram-core/resources/telegram/urls/{url_key}"
    try:
        content, content_type, final_url = await _download_public_url(url)
        FILE_ROOT.mkdir(parents=True, exist_ok=True)
        suffix = Path(urlparse(final_url).path).suffix.lower()[:16]
        if not suffix:
            suffix = (mimetypes.guess_extension(content_type) or ".bin")[:16]
        archive_path = FILE_ROOT / f"url-{url_key}{suffix}"
        if not archive_path.exists():
            archive_path.write_bytes(content)
            os.chmod(archive_path, 0o600)
        if content_type in {"text/html", "application/xhtml+xml"}:
            extracted = await asyncio.to_thread(_html_text, content)
        elif content_type.startswith("text/") or content_type in {"application/json", "application/xml"}:
            extracted = content.decode("utf-8", errors="replace")
        else:
            extracted = await asyncio.to_thread(_extract_path, archive_path, content_type)
        text = (
            f"Source URL: {url}\nFinal URL: {final_url}\nContent type: {content_type}\n"
            f"Local archive: {archive_path.name}\n\n{extracted}"
        )
    except Exception as error:
        text = f"Source URL: {url}\nFetch status: failed\nReason: {type(error).__name__}: {error}"
    await _write_chunks(base_uri, "Automatically indexed URL", text)


def _same_site_document_links(content: bytes, base_url: str) -> set[str]:
    from bs4 import BeautifulSoup

    hostname = (urlparse(base_url).hostname or "").lower().removeprefix("www.")
    links: set[str] = set()
    soup = BeautifulSoup(content, "html.parser")
    for anchor in soup.find_all("a", href=True):
        candidate = urljoin(base_url, str(anchor.get("href") or ""))
        parsed = urlparse(candidate)
        candidate_host = (parsed.hostname or "").lower().removeprefix("www.")
        if candidate_host == hostname and parsed.path.lower().endswith(".pdf"):
            links.add(candidate)
    return links


async def _discover_site_pdfs(url: str) -> list[str]:
    """Discover same-site PDFs, preferring the public WordPress media API."""
    content, content_type, final_url = await _download_public_url(url)
    parsed = urlparse(final_url)
    if content_type not in {"text/html", "application/xhtml+xml"} or not parsed.hostname:
        return []

    links = _same_site_document_links(content, final_url)
    api_url = f"{parsed.scheme}://{parsed.netloc}/wp-json/wp/v2/media?per_page=100&page=1"
    try:
        api_content, api_type, _api_final = await _download_public_url(api_url)
        if api_type == "application/json" or api_type.startswith("text/"):
            media = json.loads(api_content.decode("utf-8", errors="replace"))
            if isinstance(media, list):
                site_host = parsed.hostname.lower().removeprefix("www.")
                for item in media:
                    if not isinstance(item, dict) or item.get("mime_type") != "application/pdf":
                        continue
                    source = str(item.get("source_url") or "")
                    source_host = (urlparse(source).hostname or "").lower().removeprefix("www.")
                    if source.startswith(("http://", "https://")) and source_host == site_host:
                        links.add(source)
    except Exception:
        pass
    return sorted(links)[:MAX_SITE_DOCUMENTS]


async def _index_site_pdfs(url: str) -> None:
    try:
        document_urls = await _discover_site_pdfs(url)
    except Exception as error:
        failure_uri = f"viking://user/telegram-core/resources/telegram/sites/{_sha(url)[:32]}.md"
        await _write_content(
            failure_uri,
            f"# Site PDF discovery\n\nSource URL: {url}\nStatus: failed\n"
            f"Reason: {type(error).__name__}: {error}\n",
        )
        return

    manifest_uri = f"viking://user/telegram-core/resources/telegram/sites/{_sha(url)[:32]}.md"
    manifest = "\n".join(f"- {document_url}" for document_url in document_urls) or "[No PDFs found]"
    await _write_content(
        manifest_uri,
        f"# Site PDF discovery\n\nSource URL: {url}\nPDF count: {len(document_urls)}\n\n{manifest}\n",
    )
    for document_url in document_urls:
        await _index_url(document_url)


async def _index_archive(path: Path) -> None:
    """Submit a ZIP once to OpenViking's local archive parser."""
    import httpx

    endpoint = os.environ.get("OPENVIKING_ENDPOINT", "http://openviking:1933").rstrip("/")
    file_digest = await asyncio.to_thread(_file_sha256, path)
    target_uri = f"viking://resources/file-{file_digest[:20]}"
    headers = _headers()
    timeout = httpx.Timeout(600.0, connect=20.0)

    async with httpx.AsyncClient(timeout=timeout) as client:
        existing = await client.get(
            f"{endpoint}/api/v1/fs/stat",
            headers=headers,
            params={"uri": target_uri},
        )
        if existing.status_code < 400:
            return
        if existing.status_code not in {400, 404}:
            existing.raise_for_status()

        upload_headers = dict(headers)
        upload_headers.pop("Content-Type", None)
        with path.open("rb") as handle:
            upload = await client.post(
                f"{endpoint}/api/v1/resources/temp_upload",
                headers=upload_headers,
                files={"file": (path.name, handle, "application/zip")},
            )
        upload.raise_for_status()
        temp_file_id = upload.json().get("result", {}).get("temp_file_id", "")
        if not temp_file_id:
            raise RuntimeError("OpenViking archive upload did not return a temporary file ID")
        try:
            response = await client.post(
                f"{endpoint}/api/v1/resources",
                headers=headers,
                json={
                    "temp_file_id": temp_file_id,
                    "source_name": path.name,
                    "to": target_uri,
                    "wait": False,
                },
            )
            response.raise_for_status()
        except (httpx.TimeoutException, httpx.HTTPStatusError):
            accepted = await client.get(
                f"{endpoint}/api/v1/fs/stat",
                headers=headers,
                params={"uri": target_uri},
            )
            if accepted.status_code >= 400:
                raise


async def _index_media(item: dict[str, str], message_key: str, index: int) -> None:
    path_value = item.get("path") or ""
    source = item.get("source") or ""
    mime_type = item.get("mime") or mimetypes.guess_type(path_value)[0] or "application/octet-stream"
    base_uri = f"viking://user/telegram-core/resources/telegram/files/{message_key}-{index:02d}"
    if path_value and Path(path_value).is_file():
        path = Path(path_value)
        if path.suffix.lower() == ".zip" or mime_type == "application/zip":
            await _index_archive(path)
            return
        extracted = await asyncio.to_thread(_extract_path, path, mime_type)
        text = f"Archived file: {path.name}\nOriginal source: {source}\nMIME: {mime_type}\n\n{extracted}"
    elif source.startswith(("http://", "https://")):
        text = f"Remote attachment source: {source}\nMIME: {mime_type}"
    else:
        text = f"Attachment metadata only. Source: {source}\nMIME: {mime_type}"
    await _write_chunks(base_uri, "Automatically indexed Telegram attachment", text)


async def _process_spool(path: Path) -> None:
    payload = json.loads(path.read_text(encoding="utf-8"))
    chat_key = _sha(payload.get("chat_id", ""))[:16]
    date_key = str(payload.get("timestamp") or "unknown")[:10]
    message_uri = (
        "viking://user/telegram-core/resources/telegram/messages/"
        f"{chat_key}/{_safe_segment(date_key, 'unknown-date')}/{payload['key']}.md"
    )
    await _write_content(message_uri, _message_markdown(payload))
    for url in payload.get("urls", []):
        await _index_url(str(url))
        if payload.get("crawl_site_documents"):
            await _index_site_pdfs(str(url))
    for index, item in enumerate(payload.get("media", [])):
        await _index_media(item, payload["key"], index)
    path.unlink(missing_ok=True)


async def _drain_worker() -> None:
    """Drain the durable spool sequentially without blocking message routing."""
    assert _WORKER_WAKEUP is not None
    while True:
        pending = sorted(SPOOL_ROOT.glob("*.json"))
        if not pending:
            _WORKER_WAKEUP.clear()
            if any(SPOOL_ROOT.glob("*.json")):
                continue
            await _WORKER_WAKEUP.wait()
            continue
        for candidate in pending:
            try:
                await _process_spool(candidate)
            except Exception as error:
                print(
                    f"[tietopolitiikka-ingest] Deferred {candidate.name}: {type(error).__name__}",
                    flush=True,
                )
                await asyncio.sleep(2)


async def archive_telegram_event(event: Any) -> None:
    """Durably spool an allowed Telegram event and schedule local indexing."""
    global _WORKER_TASK, _WORKER_WAKEUP
    if not is_allowed_telegram_event(event):
        return
    try:
        data = _telegram_event_data(event)
        await asyncio.to_thread(_create_spool, event, data)
        if _WORKER_WAKEUP is None:
            _WORKER_WAKEUP = asyncio.Event()
        if _WORKER_TASK is None or _WORKER_TASK.done():
            _WORKER_TASK = asyncio.create_task(_drain_worker(), name="tietopolitiikka-ingest")
        _WORKER_WAKEUP.set()
    except Exception as error:
        print(f"[tietopolitiikka-ingest] Spool failure: {type(error).__name__}", flush=True)
